#!/usr/bin/env python3
"""
Script to create a professional PowerPoint presentation for Doctor-AI project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_doctor_ai_presentation():
    """Create a 5-slide professional presentation for Doctor-AI"""

    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Define color scheme (Medical/Healthcare theme)
    PRIMARY_COLOR = RGBColor(0, 114, 188)  # Medical blue
    SECONDARY_COLOR = RGBColor(0, 163, 224)  # Light blue
    ACCENT_COLOR = RGBColor(76, 175, 80)  # Green for positive impact
    TEXT_COLOR = RGBColor(51, 51, 51)  # Dark gray

    # SLIDE 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background color
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR

    # Title
    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Doctor-AI"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(72)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "AI-Powered Clinical Decision Support System"
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.font.size = Pt(32)
    subtitle_p.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_p.alignment = PP_ALIGN.CENTER

    # Footer
    footer_box = slide1.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Transforming Healthcare with Advanced AI Technology"
    footer_p = footer_frame.paragraphs[0]
    footer_p.font.size = Pt(18)
    footer_p.font.italic = True
    footer_p.font.color.rgb = SECONDARY_COLOR
    footer_p.alignment = PP_ALIGN.CENTER

    # SLIDE 2: Overview & Problem Statement
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame2 = title_box2.text_frame
    title_frame2.text = "What is Doctor-AI?"
    title_p2 = title_frame2.paragraphs[0]
    title_p2.font.size = Pt(44)
    title_p2.font.bold = True
    title_p2.font.color.rgb = PRIMARY_COLOR

    # Content
    content_box2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(5.5))
    tf2 = content_box2.text_frame
    tf2.word_wrap = True

    # Problem statement
    p1 = tf2.paragraphs[0]
    p1.text = "The Challenge"
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_COLOR
    p1.space_after = Pt(10)

    # Problem details
    p2 = tf2.add_paragraph()
    p2.text = "Healthcare professionals face complex diagnostic decisions with limited time and vast medical knowledge to process."
    p2.font.size = Pt(20)
    p2.font.color.rgb = TEXT_COLOR
    p2.space_after = Pt(20)
    p2.level = 0

    # Solution
    p3 = tf2.add_paragraph()
    p3.text = "The Solution"
    p3.font.size = Pt(28)
    p3.font.bold = True
    p3.font.color.rgb = ACCENT_COLOR
    p3.space_after = Pt(10)

    # Solution details
    solutions = [
        "Analyzes patient symptoms using advanced AI and semantic search",
        "Provides differential diagnoses ranked by confidence levels",
        "Detects rare diseases through comprehensive medical ontologies",
        "Flags life-threatening symptoms requiring immediate attention",
        "Offers explainable AI reasoning for all diagnostic suggestions"
    ]

    for solution in solutions:
        p = tf2.add_paragraph()
        p.text = solution
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.level = 0
        p.space_after = Pt(8)
        # Add bullet
        p.text = "• " + solution

    # SLIDE 3: Key Features
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame3 = title_box3.text_frame
    title_frame3.text = "Key Features & Capabilities"
    title_p3 = title_frame3.paragraphs[0]
    title_p3.font.size = Pt(44)
    title_p3.font.bold = True
    title_p3.font.color.rgb = PRIMARY_COLOR

    # Feature boxes
    features = [
        {
            "title": "🔬 Advanced AI Engine",
            "points": [
                "BioBERT/PubMedBERT embeddings (768-dim)",
                "Vector similarity search via Qdrant",
                "Query latency <2 seconds (p95)"
            ]
        },
        {
            "title": "🏥 Medical Intelligence",
            "points": [
                "HPO: 15,247 phenotype terms",
                "ICD-10-CM: 70,000+ diagnostic codes",
                "Rare disease detection capabilities"
            ]
        },
        {
            "title": "⚡ Smart Triage System",
            "points": [
                "4-tier confidence-based routing",
                "Red flag alert detection",
                "Specialist recommendations"
            ]
        },
        {
            "title": "🔒 Security & Compliance",
            "points": [
                "HIPAA-compliant audit logging",
                "Data encryption at rest/transit",
                "Complete audit trail"
            ]
        }
    ]

    # Create 2x2 grid of features
    start_x = 0.6
    start_y = 1.4
    box_width = 4.4
    box_height = 2.6
    gap = 0.3

    for i, feature in enumerate(features):
        row = i // 2
        col = i % 2
        x = start_x + col * (box_width + gap)
        y = start_y + row * (box_height + gap)

        # Feature box with border
        shape = slide3.shapes.add_shape(
            1,  # Rectangle
            Inches(x), Inches(y), Inches(box_width), Inches(box_height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(245, 248, 250)
        shape.line.color.rgb = SECONDARY_COLOR
        shape.line.width = Pt(2)

        # Feature title
        text_frame = shape.text_frame
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_left = Inches(0.15)
        text_frame.margin_right = Inches(0.15)
        text_frame.word_wrap = True

        p_title = text_frame.paragraphs[0]
        p_title.text = feature["title"]
        p_title.font.size = Pt(20)
        p_title.font.bold = True
        p_title.font.color.rgb = PRIMARY_COLOR
        p_title.space_after = Pt(8)

        # Feature points
        for point in feature["points"]:
            p = text_frame.add_paragraph()
            p.text = "• " + point
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(4)

    # SLIDE 4: Technology Stack & Architecture
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame4 = title_box4.text_frame
    title_frame4.text = "Technology Stack & Architecture"
    title_p4 = title_frame4.paragraphs[0]
    title_p4.font.size = Pt(44)
    title_p4.font.bold = True
    title_p4.font.color.rgb = PRIMARY_COLOR

    # Technology sections
    tech_sections = [
        {
            "title": "Backend & AI/ML",
            "items": [
                "Python 3.9+ & FastAPI",
                "BioBERT, PubMedBERT, Transformers",
                "Qdrant Vector Database",
                "LLM Integration (GPT-4/Llama2)"
            ]
        },
        {
            "title": "Frontend & Data",
            "items": [
                "React 18 with Vite",
                "Responsive modern UI",
                "PostgreSQL database",
                "Redis caching layer"
            ]
        },
        {
            "title": "Deployment & DevOps",
            "items": [
                "Docker & Docker Compose",
                "Render.com cloud ready",
                "100+ queries/second throughput",
                "Scalable microservices architecture"
            ]
        }
    ]

    y_position = 1.3
    for section in tech_sections:
        # Section box
        shape = slide4.shapes.add_shape(
            1,
            Inches(0.8), Inches(y_position), Inches(8.4), Inches(1.6)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(240, 247, 255)
        shape.line.color.rgb = SECONDARY_COLOR
        shape.line.width = Pt(1.5)

        text_frame = shape.text_frame
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_left = Inches(0.2)
        text_frame.word_wrap = True

        # Section title
        p_title = text_frame.paragraphs[0]
        p_title.text = section["title"]
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = PRIMARY_COLOR
        p_title.space_after = Pt(10)

        # Items in a row
        p_items = text_frame.add_paragraph()
        p_items.text = "  •  ".join(section["items"])
        p_items.font.size = Pt(16)
        p_items.font.color.rgb = TEXT_COLOR

        y_position += 1.8

    # SLIDE 5: Impact & Results
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame5 = title_box5.text_frame
    title_frame5.text = "Impact & Results"
    title_p5 = title_frame5.paragraphs[0]
    title_p5.font.size = Pt(44)
    title_p5.font.bold = True
    title_p5.font.color.rgb = PRIMARY_COLOR

    # Impact metrics
    metrics = [
        {"value": "<2s", "label": "Diagnostic Query Latency (p95)"},
        {"value": "100+", "label": "Queries Per Second"},
        {"value": "15K+", "label": "Phenotype Terms (HPO)"},
        {"value": "70K+", "label": "ICD-10-CM Diagnostic Codes"}
    ]

    # Create metrics boxes
    box_width = 2.1
    box_height = 1.4
    gap = 0.2
    start_x = 0.7
    start_y = 1.4

    for i, metric in enumerate(metrics):
        col = i % 4
        x = start_x + col * (box_width + gap)

        # Metric box
        shape = slide5.shapes.add_shape(
            1,
            Inches(x), Inches(start_y), Inches(box_width), Inches(box_height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT_COLOR
        shape.line.color.rgb = RGBColor(255, 255, 255)
        shape.line.width = Pt(0)

        text_frame = shape.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_frame.word_wrap = True

        # Value
        p_value = text_frame.paragraphs[0]
        p_value.text = metric["value"]
        p_value.font.size = Pt(36)
        p_value.font.bold = True
        p_value.font.color.rgb = RGBColor(255, 255, 255)
        p_value.alignment = PP_ALIGN.CENTER

        # Label
        p_label = text_frame.add_paragraph()
        p_label.text = metric["label"]
        p_label.font.size = Pt(12)
        p_label.font.color.rgb = RGBColor(255, 255, 255)
        p_label.alignment = PP_ALIGN.CENTER

    # Key Impacts section
    impacts_box = slide5.shapes.add_textbox(Inches(0.7), Inches(3.2), Inches(8.6), Inches(3.5))
    tf_impacts = impacts_box.text_frame
    tf_impacts.word_wrap = True

    # Impacts title
    p_impacts_title = tf_impacts.paragraphs[0]
    p_impacts_title.text = "Key Impacts on Healthcare"
    p_impacts_title.font.size = Pt(28)
    p_impacts_title.font.bold = True
    p_impacts_title.font.color.rgb = PRIMARY_COLOR
    p_impacts_title.space_after = Pt(15)

    # Impact points
    impacts = [
        "✓ Assists healthcare professionals with rapid, evidence-based diagnostic support",
        "✓ Identifies rare diseases that might be missed in traditional workflows",
        "✓ Reduces diagnostic uncertainty through AI-powered differential diagnosis",
        "✓ Improves patient safety with automated red flag detection",
        "✓ Provides explainable AI reasoning for transparency and trust",
        "✓ Enhances clinical decision-making while maintaining human oversight"
    ]

    for impact in impacts:
        p = tf_impacts.add_paragraph()
        p.text = impact
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(10)

    # Footer note
    footer_box5 = slide5.shapes.add_textbox(Inches(0.7), Inches(6.9), Inches(8.6), Inches(0.5))
    footer_frame5 = footer_box5.text_frame
    footer_frame5.text = "⚕️ Clinical Decision Support Tool - Designed to Assist, Not Replace Healthcare Professionals"
    footer_p5 = footer_frame5.paragraphs[0]
    footer_p5.font.size = Pt(14)
    footer_p5.font.italic = True
    footer_p5.font.color.rgb = RGBColor(100, 100, 100)
    footer_p5.alignment = PP_ALIGN.CENTER

    # Save presentation
    output_file = "Doctor-AI_Presentation.pptx"
    prs.save(output_file)
    print(f"✓ Presentation created successfully: {output_file}")
    print(f"✓ Total slides: {len(prs.slides)}")
    return output_file

if __name__ == "__main__":
    create_doctor_ai_presentation()
